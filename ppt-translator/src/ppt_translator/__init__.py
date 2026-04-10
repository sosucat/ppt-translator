from __future__ import annotations

import argparse
import difflib
import re
from pathlib import Path

from deep_translator import GoogleTranslator
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def translate_text(text: str | None, translator: GoogleTranslator) -> str:
    if not isinstance(text, str):
        return ""
    text = text.strip()
    if not text:
        return ""
    try:
        translated = translator.translate(text)
        if isinstance(translated, str):
            return translated
        if translated is None:
            return ""
        return str(translated)
    except Exception:
        return text


def normalize_text(text: str) -> str:
    return re.sub(r"\s+", " ", text.strip().lower())


def is_english_text(text: str) -> bool:
    if not text:
        return False
    if re.search(r"[\u3040-\u30ff\u3400-\u4dbf\u4e00-\u9fff\uf900-\ufaff]", text):
        return False
    return bool(re.search(r"[A-Za-z]", text))


def get_slide_english_texts(slide) -> dict[str, str]:
    english_texts: dict[str, str] = {}
    shapes_to_process = list(slide.shapes)
    while shapes_to_process:
        shape = shapes_to_process.pop()
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            shapes_to_process.extend(shape.shapes)
            continue

        if hasattr(shape, "has_table") and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    cell_text = cell.text
                    if isinstance(cell_text, str) and cell_text.strip() and is_english_text(cell_text):
                        english_texts[normalize_text(cell_text)] = cell_text

        if hasattr(shape, "has_text_frame") and shape.has_text_frame and shape.text_frame is not None:
            for paragraph in shape.text_frame.paragraphs:
                paragraph_text = paragraph.text
                if isinstance(paragraph_text, str) and paragraph_text.strip() and is_english_text(paragraph_text):
                    english_texts[normalize_text(paragraph_text)] = paragraph_text

    return english_texts


def find_best_english_phrase(translated_text: str, english_texts: dict[str, str]) -> str | None:
    if not translated_text or not english_texts:
        return None

    normalized_translated = normalize_text(translated_text)
    if normalized_translated in english_texts:
        return english_texts[normalized_translated]

    best_match: str | None = None
    best_ratio = 0.0
    for normalized_original, original_text in english_texts.items():
        ratio = difflib.SequenceMatcher(None, normalized_translated, normalized_original).ratio()
        if ratio > best_ratio:
            best_ratio = ratio
            best_match = original_text

    return best_match if best_ratio >= 0.75 else None


def translate_presentation(input_path: Path, output_path: Path | None = None) -> Path:
    presentation = Presentation(str(input_path))
    translator = GoogleTranslator(source="ja", target="en")

    for slide in presentation.slides:
        english_texts = get_slide_english_texts(slide)
        shapes_to_process = list(slide.shapes)
        while shapes_to_process:
            shape = shapes_to_process.pop()
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                shapes_to_process.extend(shape.shapes)
                continue

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                continue

            if hasattr(shape, "has_table") and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell_text = cell.text
                        if isinstance(cell_text, str) and cell_text.strip():
                            translated = translate_text(cell_text, translator)
                            best_phrase = find_best_english_phrase(translated, english_texts)
                            if best_phrase is not None:
                                cell.text = ""
                            else:
                                cell.text = translated
                continue

            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                if shape.text_frame is None:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = paragraph.text
                    if isinstance(paragraph_text, str) and paragraph_text.strip():
                        translated = translate_text(paragraph_text, translator)
                        best_phrase = find_best_english_phrase(translated, english_texts)
                        if best_phrase is not None:
                            paragraph.text = ""
                        elif translated:
                            paragraph.text = translated
                continue

    if output_path is None:
        output_path = input_path.with_name(f"{input_path.stem}_translated{input_path.suffix}")

    presentation.save(str(output_path))
    return output_path


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Translate Japanese text in a PPTX file to English.")
    parser.add_argument("input", type=Path, help="Path to the source .pptx file.")
    parser.add_argument("--output", "-o", type=Path, help="Optional output path for the translated .pptx file.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    input_path = args.input

    if not input_path.exists() or input_path.suffix.lower() != ".pptx":
        raise SystemExit("Please provide a valid .pptx file path.")

    output_path = args.output if args.output is not None else None
    result_path = translate_presentation(input_path, output_path)
    print(f"Translated presentation saved to: {result_path}")


if __name__ == "__main__":
    main()
